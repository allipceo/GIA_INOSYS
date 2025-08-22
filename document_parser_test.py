#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
GIA_INFOSYS 문서 파싱 테스트 스크립트
Phase 2-1, 과업 1: 문서 파싱 환경 구축

작성일: 2025년 8월 22일
작성자: 서대리 (Lead Developer)
목적: 다양한 문서 형식(DOCX, PPTX, PDF)에서 텍스트 추출 기능 테스트
"""

import os
import sys
from datetime import datetime
from typing import Optional, Dict, Any

# 문서 파싱 라이브러리
from docx import Document
from pptx import Presentation
import fitz  # PyMuPDF

# Notion API
from notion_client import Client
import requests

# 환경 변수
from dotenv import load_dotenv

# 로깅 설정
import logging

# 환경 변수 로드
load_dotenv()

# 로깅 설정
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('document_parser.log'),
        logging.StreamHandler(sys.stdout)
    ]
)
logger = logging.getLogger(__name__)

class DocumentParser:
    """문서 파싱 클래스"""
    
    def __init__(self):
        """초기화"""
        self.notion_token = os.getenv('NOTION_TOKEN')
        self.notion_database_id = os.getenv('NOTION_DATABASE_ID')
        
        if self.notion_token:
            self.notion_client = Client(auth=self.notion_token)
        else:
            self.notion_client = None
            logger.warning("Notion 토큰이 설정되지 않았습니다.")
    
    def parse_docx(self, file_path: str) -> Optional[str]:
        """
        .docx 파일에서 텍스트를 추출하는 함수
        
        Args:
            file_path (str): DOCX 파일 경로
            
        Returns:
            str: 추출된 텍스트
        """
        try:
            if not os.path.exists(file_path):
                logger.error(f"파일이 존재하지 않습니다: {file_path}")
                return None
            
            doc = Document(file_path)
            text_content = []
            
            # 단락에서 텍스트 추출
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text)
            
            # 테이블에서 텍스트 추출
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text_content.append(" | ".join(row_text))
            
            extracted_text = "\n".join(text_content)
            logger.info(f"DOCX 파일 파싱 성공: {file_path}")
            return extracted_text
            
        except Exception as e:
            logger.error(f"DOCX 파일 파싱 실패: {file_path}, 오류: {str(e)}")
            return None
    
    def parse_pptx(self, file_path: str) -> Optional[str]:
        """
        .pptx 파일에서 텍스트를 추출하는 함수
        
        Args:
            file_path (str): PPTX 파일 경로
            
        Returns:
            str: 추출된 텍스트
        """
        try:
            if not os.path.exists(file_path):
                logger.error(f"파일이 존재하지 않습니다: {file_path}")
                return None
            
            prs = Presentation(file_path)
            text_content = []
            
            # 각 슬라이드에서 텍스트 추출
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                
                # 슬라이드 제목
                if slide.shapes.title:
                    slide_text.append(f"슬라이드 {slide_num} 제목: {slide.shapes.title.text}")
                
                # 슬라이드 내용
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if slide_text:
                    text_content.append("\n".join(slide_text))
            
            extracted_text = "\n\n".join(text_content)
            logger.info(f"PPTX 파일 파싱 성공: {file_path}")
            return extracted_text
            
        except Exception as e:
            logger.error(f"PPTX 파일 파싱 실패: {file_path}, 오류: {str(e)}")
            return None
    
    def parse_pdf(self, file_path: str) -> Optional[str]:
        """
        .pdf 파일에서 텍스트를 추출하는 함수
        
        Args:
            file_path (str): PDF 파일 경로
            
        Returns:
            str: 추출된 텍스트
        """
        try:
            if not os.path.exists(file_path):
                logger.error(f"파일이 존재하지 않습니다: {file_path}")
                return None
            
            doc = fitz.open(file_path)
            text_content = []
            
            # 각 페이지에서 텍스트 추출
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                page_text = page.get_text()
                if page_text.strip():
                    text_content.append(f"페이지 {page_num + 1}:\n{page_text.strip()}")
            
            doc.close()
            extracted_text = "\n\n".join(text_content)
            logger.info(f"PDF 파일 파싱 성공: {file_path}")
            return extracted_text
            
        except Exception as e:
            logger.error(f"PDF 파일 파싱 실패: {file_path}, 오류: {str(e)}")
            return None
    
    def create_test_files(self) -> Dict[str, str]:
        """
        테스트용 샘플 파일들을 생성하는 함수
        
        Returns:
            Dict[str, str]: 생성된 파일 경로들
        """
        test_files = {}
        
        # 테스트 파일 디렉토리 생성
        test_dir = "./test_files"
        os.makedirs(test_dir, exist_ok=True)
        
        # DOCX 테스트 파일 생성
        docx_path = os.path.join(test_dir, "test.docx")
        doc = Document()
        doc.add_heading('GIA_INFOSYS 테스트 문서', 0)
        doc.add_paragraph('이것은 DOCX 파일 파싱 테스트를 위한 샘플 문서입니다.')
        doc.add_paragraph('다양한 내용을 포함하고 있습니다.')
        doc.add_paragraph('테스트 성공 시 Notion DB에 저장됩니다.')
        doc.save(docx_path)
        test_files['docx'] = docx_path
        
        # PPTX 테스트 파일 생성
        pptx_path = os.path.join(test_dir, "test.pptx")
        prs = Presentation()
        slide_layout = prs.slide_layouts[0]  # 제목 슬라이드
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "GIA_INFOSYS 테스트 프레젠테이션"
        slide.placeholders[1].text = "이것은 PPTX 파일 파싱 테스트를 위한 샘플 프레젠테이션입니다.\n\n다양한 슬라이드 내용을 포함합니다.\n\n테스트 성공 시 Notion DB에 저장됩니다."
        prs.save(pptx_path)
        test_files['pptx'] = pptx_path
        
        # PDF 테스트 파일 생성 (텍스트 기반)
        pdf_path = os.path.join(test_dir, "test.pdf")
        doc = fitz.open()
        page = doc.new_page()
        text = "GIA_INFOSYS 테스트 PDF\n\n이것은 PDF 파일 파싱 테스트를 위한 샘플 문서입니다.\n\n다양한 내용을 포함하고 있습니다.\n\n테스트 성공 시 Notion DB에 저장됩니다."
        page.insert_text((50, 50), text)
        doc.save(pdf_path)
        doc.close()
        test_files['pdf'] = pdf_path
        
        logger.info("테스트 파일 생성 완료")
        return test_files
    
    def test_all_parsers(self) -> Dict[str, Any]:
        """
        모든 파서를 테스트하는 함수
        
        Returns:
            Dict[str, Any]: 테스트 결과
        """
        logger.info("=== 문서 파싱 테스트 시작 ===")
        
        # 테스트 파일 생성
        test_files = self.create_test_files()
        results = {}
        
        # 각 파일 형식별 파싱 테스트
        for file_type, file_path in test_files.items():
            logger.info(f"테스트 중: {file_type.upper()} 파일")
            
            if file_type == 'docx':
                result = self.parse_docx(file_path)
            elif file_type == 'pptx':
                result = self.parse_pptx(file_path)
            elif file_type == 'pdf':
                result = self.parse_pdf(file_path)
            else:
                continue
            
            results[file_type] = {
                'file_path': file_path,
                'success': result is not None,
                'text_length': len(result) if result else 0,
                'extracted_text': result[:200] + "..." if result and len(result) > 200 else result
            }
            
            if result:
                logger.info(f"✅ {file_type.upper()} 파싱 성공: {len(result)} 문자")
            else:
                logger.error(f"❌ {file_type.upper()} 파싱 실패")
        
        logger.info("=== 문서 파싱 테스트 완료 ===")
        return results
    
    def create_notion_database(self) -> Optional[str]:
        """
        Notion에 임시 테스트용 DB를 생성하는 함수
        
        Returns:
            str: 생성된 DB ID
        """
        if not self.notion_client:
            logger.error("Notion 클라이언트가 초기화되지 않았습니다.")
            return None
        
        try:
            # DB 생성 (실제로는 수동으로 생성하고 ID만 반환)
            logger.info("Notion DB 생성 기능은 수동으로 수행해야 합니다.")
            logger.info("MainGate 페이지에 '임시 테스트용 DB'를 생성하고 DB ID를 .env 파일에 설정하세요.")
            return None
            
        except Exception as e:
            logger.error(f"Notion DB 생성 실패: {str(e)}")
            return None

def main():
    """메인 실행 함수"""
    print("=== GIA_INFOSYS 문서 파싱 테스트 시작 ===")
    
    # DocumentParser 인스턴스 생성
    parser = DocumentParser()
    
    # 모든 파서 테스트 실행
    results = parser.test_all_parsers()
    
    # 결과 출력
    print("\n=== 테스트 결과 ===")
    for file_type, result in results.items():
        status = "✅ 성공" if result['success'] else "❌ 실패"
        print(f"{file_type.upper()}: {status}")
        if result['success']:
            print(f"  - 파일: {result['file_path']}")
            print(f"  - 텍스트 길이: {result['text_length']} 문자")
            print(f"  - 추출된 텍스트 미리보기: {result['extracted_text']}")
        print()
    
    # 전체 성공 여부 확인
    all_success = all(result['success'] for result in results.values())
    if all_success:
        print("🎉 모든 문서 파싱 테스트가 성공했습니다!")
    else:
        print("⚠️ 일부 문서 파싱 테스트가 실패했습니다.")
    
    print("\n=== 다음 단계 ===")
    print("1. Notion DB 생성 (수동)")
    print("2. .env 파일에 Notion 토큰과 DB ID 설정")
    print("3. Notion 연동 테스트 실행")
    
    return all_success

if __name__ == "__main__":
    main()
